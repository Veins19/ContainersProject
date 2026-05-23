from fastapi import FastAPI, Request, WebSocket
import chromadb
from chromadb import Documents, EmbeddingFunction, Embeddings
import google.generativeai as genai
import uuid
import io
import pypdf
import docx
from pptx import Presentation
import asyncio
import os
import requests
import hashlib
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import CrossEncoder

load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "YOUR_PLACEHOLDER_KEY_HERE")
genai.configure(api_key=GEMINI_API_KEY)

gemini_model = genai.GenerativeModel('gemini-3.1-flash-lite')

re_ranker = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')

class GeminiEmbeddingFunction(EmbeddingFunction):
    def __call__(self, input: Documents) -> Embeddings:
        response = genai.embed_content(
            model="models/gemini-embedding-001",
            content=input
        )
        return response['embedding']

app = FastAPI()

chroma_client = chromadb.PersistentClient(path="./chroma_storage")

collection = chroma_client.get_or_create_collection(
    name="lab_knowledge",
    embedding_function=GeminiEmbeddingFunction()
)

@app.websocket("/ws/upload")
async def websocket_upload(websocket: WebSocket):
    await websocket.accept()
    try:
        meta = await websocket.receive_json()
        ext = meta.get("ext", "").lower()
        user_id = meta.get("user_id")
        document_id = meta.get("document_id")
        
        if not user_id or not document_id:
            await websocket.send_json({"progress": 100, "status": "Error", "message": "Missing IDs."})
            await websocket.close()
            return

        content = await websocket.receive_bytes()
        text = ""
        
        if ext == "pdf":
            pdf_reader = pypdf.PdfReader(io.BytesIO(content))
            total_pages = len(pdf_reader.pages)
            for i, page in enumerate(pdf_reader.pages):
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
                progress = int(((i + 1) / total_pages) * 90)
                await websocket.send_json({"progress": progress, "status": f"Reading PDF page {i+1}..."})
                await asyncio.sleep(0.01) 
        elif ext == "docx":
            doc = docx.Document(io.BytesIO(content))
            total_paras = len(doc.paragraphs)
            for i, para in enumerate(doc.paragraphs):
                text += para.text + "\n"
                if i % max(1, total_paras // 20) == 0:
                    progress = int(((i + 1) / total_paras) * 90)
                    await websocket.send_json({"progress": progress, "status": "Parsing Word..."})
                    await asyncio.sleep(0.01)
        elif ext == "pptx":
            # THE PPTX UPGRADE: Slide-Aware Anchoring
            ppt = Presentation(io.BytesIO(content))
            total_slides = len(ppt.slides)
            for i, slide in enumerate(ppt.slides):
                slide_text = f"--- Slide {i+1} ---\n"
                has_content = False
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Clean weird PowerPoint invisible characters
                        clean_text = shape.text.replace('\x0b', '\n').strip()
                        slide_text += clean_text + "\n"
                        has_content = True
                
                # Only add the slide if it actually had text
                if has_content:
                    # The double newline guarantees the semantic chunker respects the slide boundary
                    text += slide_text + "\n\n" 
                
                progress = int(((i + 1) / total_slides) * 90)
                await websocket.send_json({"progress": progress, "status": f"Reading Slide {i+1}..."})
                await asyncio.sleep(0.01)
        elif ext == "txt":
            text = content.decode('utf-8')
            await websocket.send_json({"progress": 90, "status": "Reading text..."})

        await websocket.send_json({"progress": 95, "status": "Hashing and Syncing with Vector DB..."})
        
        if text.strip():
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1200,
                chunk_overlap=250,
                separators=["\n\n", "\n", ".", " ", ""]
            )
            chunks = text_splitter.split_text(text)
            
            new_chunk_ids = []
            new_chunks = []
            new_metadatas = []
            
            for chunk in chunks:
                chunk_hash = hashlib.sha256(chunk.encode('utf-8')).hexdigest()
                doc_chunk_id = f"{document_id}_{chunk_hash}"
                new_chunk_ids.append(doc_chunk_id)
                new_chunks.append(chunk)
                new_metadatas.append({"user_id": user_id, "document_id": document_id, "hash": chunk_hash})
            
            existing_docs = collection.get(
                where={
                    "$and": [
                        {"user_id": user_id},
                        {"document_id": document_id}
                    ]
                }
            )
            existing_ids = set(existing_docs['ids']) if existing_docs and existing_docs['ids'] else set()
            new_ids_set = set(new_chunk_ids)
            
            ids_to_delete = list(existing_ids - new_ids_set)
            if ids_to_delete:
                collection.delete(ids=ids_to_delete)
                
            if new_chunk_ids:
                collection.upsert(
                    documents=new_chunks, 
                    metadatas=new_metadatas, 
                    ids=new_chunk_ids
                )
                
            await websocket.send_json({"progress": 100, "status": "Success", "message": "Successfully synchronized secure document!"})
        else:
            await websocket.send_json({"progress": 100, "status": "Error", "message": "No extractable text found."})

    except Exception as e:
        await websocket.send_json({"progress": 0, "status": "Error", "message": f"Parsing Failed: {str(e)}"})
    finally:
        await websocket.close()

@app.post("/ask")
async def ask_ai(request: Request):
    data = await request.json()
    user_question = data.get("question", "")
    user_id = data.get("user_id")
    engine = data.get("engine", "gemini")
    chat_history = data.get("chat_history", [])

    if not user_id:
         return {"status": "Error", "answer": "Unauthorized"}

    retrieved_context = ""
    
    try:
        if collection.count() > 0:
            results = collection.query(
                query_texts=[user_question], 
                n_results=15, 
                where={"user_id": user_id} 
            )
            
            if results['documents'] and len(results['documents'][0]) > 0:
                raw_chunks = results['documents'][0]
                cross_inp = [[user_question, chunk] for chunk in raw_chunks]
                cross_scores = re_ranker.predict(cross_inp)
                scored_chunks = sorted(zip(cross_scores, raw_chunks), reverse=True)
                top_3_chunks = [chunk for score, chunk in scored_chunks[:3]]
                retrieved_context = "\n\n---\n\n".join(top_3_chunks)
    except Exception as db_error:
        return {"status": "Error", "answer": f"DATABASE RETRIEVAL ERROR: {str(db_error)}"}
        
    history_text = ""
    for turn in chat_history:
        history_text += f"User: {turn.get('user')}\nAI: {turn.get('assistant')}\n\n"
        
    prompt = f"""
    You are the RAG Neural Engine, a highly intelligent and secure AI assistant.
    
    Context from the user's secure vault: 
    "{retrieved_context}"
    
    Recent Conversation History:
    {history_text if history_text else "No prior history."}
    
    Instructions:
    1. Read the Conversation History to understand pronouns or references like "that" or "this".
    2. If the user says hello or makes small talk, respond naturally.
    3. If the user asks a factual question, answer ONLY using the Context.
    4. If the Context is empty or irrelevant, state you cannot find that info in their documents.
    
    User's Current Input: {user_question}
    """
    
    display_context = retrieved_context if retrieved_context else "No context needed."

    if engine == "ollama":
        possible_hosts = ["host.docker.internal", "192.168.65.2", "172.18.0.1", "172.17.0.1", "10.0.2.2"]
        ollama_url = None

        for host in possible_hosts:
            try:
                requests.get(f"http://{host}:11434/", timeout=1)
                ollama_url = f"http://{host}:11434/api/generate"
                break
            except requests.exceptions.RequestException:
                continue

        if not ollama_url:
            return {"status": "Error", "answer": "OLLAMA CONNECTION ERROR: Make sure Ollama is running in PowerShell with OLLAMA_HOST='0.0.0.0'"}

        try:
            payload = {
                "model": "gemma4:e2b", 
                "prompt": prompt,
                "stream": False
            }
            ollama_response = requests.post(ollama_url, json=payload, timeout=600)
            ollama_response.raise_for_status()
            answer = ollama_response.json().get("response", "Error generating local response.")
            return {"status": "Success", "retrieved_context": display_context, "answer": answer}
        except requests.exceptions.RequestException as e:
             return {"status": "Error", "answer": f"OLLAMA TIMEOUT Details: {str(e)}"}
    else:
        try:
            response = gemini_model.generate_content(prompt)
            return {"status": "Success", "retrieved_context": display_context, "answer": response.text}
        except Exception as e:
            return {"status": "Error", "answer": f"GEMINI API ERROR: {str(e)}"}

@app.delete("/documents/{document_id}")
async def delete_document(document_id: str, request: Request):
    data = await request.json()
    user_id = data.get("user_id")
    
    if not user_id:
        return {"status": "Error", "message": "Unauthorized"}
        
    try:
        collection.delete(
            where={
                "$and": [
                    {"user_id": user_id},
                    {"document_id": document_id}
                ]
            }
        )
        return {"status": "Success", "message": "Memory purged."}
    except Exception as e:
        return {"status": "Error", "message": str(e)}
